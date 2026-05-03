from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_PATH = r"C:\Users\kprv9\.gemini\antigravity\scratch\CyberShield_AI_Algorithms_Presentation.pptx"

BG = RGBColor(246, 248, 251)
NAVY = RGBColor(20, 35, 60)
CYAN = RGBColor(13, 148, 136)
ORANGE = RGBColor(234, 88, 12)
TEXT = RGBColor(37, 43, 54)
MUTED = RGBColor(96, 107, 122)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(226, 232, 240)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title, subtitle=""):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(10.5), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.47), Inches(0.48), Inches(11.3), Inches(0.2))
        p2 = sb.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = RGBColor(214, 225, 240)


def add_footer(slide, text="CyberShield AI"):
    box = slide.shapes.add_textbox(Inches(10.8), Inches(7.02), Inches(2.0), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_bullets(slide, items, left, top, width, height, size=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(7)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size if level == 0 else size - 2)
            run.font.color.rgb = TEXT if level == 0 else MUTED
            run.font.bold = level == 0 and len(text) < 70


def add_two_col(slide, left_title, left_items, right_title, right_items):
    for x, color, title in [(0.6, CYAN, left_title), (6.95, ORANGE, right_title)]:
        hdr = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.1), Inches(5.75), Inches(0.42)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        hdr.text_frame.text = title
        run = hdr.text_frame.paragraphs[0].runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = WHITE

    add_bullets(slide, left_items, 0.72, 1.66, 5.45, 5.4, 18)
    add_bullets(slide, right_items, 7.05, 1.66, 5.45, 5.4, 18)


def add_algo_slide(prs, title, what_it_does, used_for, examples):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, "Algorithm purpose, usage, and practical role")

    card1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.2), Inches(12.0), Inches(1.3)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = LIGHT
    add_bullets(slide, ["What It Does", (what_it_does, 1)], 0.9, 1.35, 11.3, 0.95, 18)

    card2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.75), Inches(5.8), Inches(3.55)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = LIGHT
    add_bullets(slide, ["Used For"] + [(item, 1) for item in used_for], 0.9, 2.95, 5.15, 3.05, 17)

    card3 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.75), Inches(5.85), Inches(3.55)
    )
    card3.fill.solid()
    card3.fill.fore_color.rgb = WHITE
    card3.line.color.rgb = LIGHT
    add_bullets(slide, ["Where It Appears In The Project"] + [(item, 1) for item in examples], 7.05, 2.95, 5.15, 3.05, 17)
    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.55)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(11.0), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "CyberShield AI"
    r.font.name = "Aptos Display"
    r.font.size = Pt(29)
    r.font.bold = True
    r.font.color.rgb = NAVY
    p2 = t.text_frame.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Algorithms Used And What They Do"
    r2.font.name = "Aptos"
    r2.font.size = Pt(22)
    r2.font.color.rgb = CYAN

    add_bullets(slide, [
        "This presentation explains the main algorithms used in the project, how they work, and why they are important.",
        "The project combines heuristics, scoring models, lightweight NLP, hashing, confidence ranking, and distance calculation.",
    ], 0.82, 2.25, 11.6, 1.3, 20)

    add_two_col(
        slide,
        "Why Algorithms Matter",
        [
            "They turn raw user input into meaningful security decisions.",
            "They help classify content as safe, suspicious, or malicious.",
            "They improve consistency instead of relying on guesswork.",
        ],
        "Project Areas",
        [
            "Scam detection",
            "Password strength",
            "File scanning",
            "URL safety",
            "App privacy risk",
            "Transaction fraud analysis",
            "Footprint confidence",
        ],
    )
    add_footer(slide)

    add_algo_slide(
        prs,
        "1. Rule-Based Heuristic Analysis",
        "Uses predefined rules, regex patterns, suspicious keywords, and safe-pattern checks to detect risk signals quickly.",
        [
            "Phishing text analysis",
            "URL and domain risk detection",
            "Suspicious file-name and extension checks",
            "Transaction fraud indicators",
        ],
        [
            "Scam detection text-signal engine",
            "Safe browsing checker",
            "File scanner local pre-scan",
            "Transaction analyzer flag generation",
        ],
    )

    add_algo_slide(
        prs,
        "2. Weighted Scoring Algorithm",
        "Assigns positive and negative points to signals and converts them into a final score or risk level.",
        [
            "Password strength scoring",
            "Scam threat scoring",
            "App privacy scoring",
            "Transaction risk scoring",
            "Footprint confidence scoring",
        ],
        [
            "Password checker strength model",
            "Threat score in scam detector",
            "App risk aggregation model",
            "Transaction analyzer score",
        ],
    )

    add_algo_slide(
        prs,
        "3. Lexicon-Based NLP Classifier",
        "Splits text into tokens, compares them with suspicious and legitimate word dictionaries, and builds a phishing-oriented language score.",
        [
            "Detecting phishing wording",
            "Estimating how scam-like a message sounds",
            "Reducing false positives with legitimate vocabulary",
        ],
        [
            "Tokenization of phishing text",
            "Suspicious and legitimate lexicon weights",
            "Browser NLP model used in AI Scam Detection",
        ],
    )

    add_algo_slide(
        prs,
        "4. Sigmoid Probability Conversion",
        "Transforms a raw NLP score into a bounded probability value so the output is easier to read as confidence.",
        [
            "Phishing probability output",
            "Confidence-style message classification",
        ],
        [
            "Used after NLP raw-score calculation",
            "Produces probability-style phishing estimate",
        ],
    )

    add_algo_slide(
        prs,
        "5. SHA-256 Hashing",
        "Creates a unique cryptographic fingerprint for a file so it can be identified without relying only on the file name.",
        [
            "File identity checking",
            "VirusTotal hash-based report lookup",
            "Safer malware analysis workflow",
        ],
        [
            "File scanner hashing step",
            "Local hash generation before optional API lookup",
        ],
    )

    add_algo_slide(
        prs,
        "6. Domain And URL Risk Analysis",
        "Examines domain structure for phishing traits such as suspicious TLDs, IP-based URLs, excessive subdomains, punycode, and deceptive trust words.",
        [
            "Website safety checks",
            "Suspicious link inspection in scam detection",
            "Brand-domain mismatch detection",
        ],
        [
            "Safe browsing checker",
            "Scam detection URL analyzer",
            "Live landing-page inspection logic",
        ],
    )

    add_algo_slide(
        prs,
        "7. Permission And Tracker Risk Aggregation",
        "Combines permission sensitivity, tracker presence, network signals, engine verdicts, and publisher trust into a privacy-risk score.",
        [
            "App privacy scanning",
            "Permission abuse estimation",
            "Tracker-heavy app identification",
        ],
        [
            "App privacy scanner score model",
            "Vendor-matrix style reputation aggregation",
            "Unknown-app fallback risk model",
        ],
    )

    add_algo_slide(
        prs,
        "8. Behavioral Transaction Risk Model",
        "Evaluates money movement using context such as amount anomaly, recipient trust, device trust, location trust, scenario, and user history.",
        [
            "Payment fraud detection",
            "Suspicious transaction review",
            "Decision support before sending money",
        ],
        [
            "Transaction analyzer flags",
            "History-based recipient and amount checks",
            "Scenario-specific fraud scoring",
        ],
    )

    add_algo_slide(
        prs,
        "9. Confidence Ranking",
        "Assigns confidence labels and sorts results so the strongest evidence appears first instead of mixing guesses with verified signals.",
        [
            "Footprint tracker result ordering",
            "Public-trace evidence explanation",
            "Reducing misleading search output",
        ],
        [
            "Result-confidence helper functions",
            "Sorted manual and verified results",
            "Accuracy summaries in footprint reports",
        ],
    )

    add_algo_slide(
        prs,
        "10. Haversine Distance Formula",
        "Calculates the distance between two GPS coordinate points on Earth using latitude and longitude values.",
        [
            "Nearest police station guidance",
            "Route and location calculations",
            "Live GPS distance estimation",
        ],
        [
            "GPS tracker module",
            "Nearby location and route support",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "Conclusion", "Summary of algorithmic design")
    add_bullets(slide, [
        "The project does not rely on one single algorithm. It uses a layered approach.",
        "Heuristics provide quick first-level detection.",
        "Weighted scoring turns many small signals into practical risk outputs.",
        "NLP improves scam-message understanding.",
        "Hashing and domain analysis improve technical accuracy.",
        "Confidence ranking helps present honest and understandable results.",
        "Together, these algorithms make the dashboard practical for real-world cybersecurity support."
    ], 0.85, 1.4, 11.8, 5.5, 20)
    add_footer(slide, "Algorithms Presentation")

    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
