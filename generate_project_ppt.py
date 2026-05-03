from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_PATH = r"C:\Users\kprv9\.gemini\antigravity\scratch\CyberShield_AI_Project_Presentation.pptx"


BG = RGBColor(247, 249, 252)
NAVY = RGBColor(16, 33, 62)
CYAN = RGBColor(17, 138, 178)
ACCENT = RGBColor(237, 108, 2)
TEXT = RGBColor(36, 41, 46)
MUTED = RGBColor(92, 102, 117)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.15), Inches(9.6), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.47), Inches(0.5), Inches(10.8), Inches(0.22))
        p2 = sub_box.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "Aptos"
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = RGBColor(213, 225, 240)


def add_bullets(slide, items, left=0.7, top=1.2, width=12.0, height=5.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for index, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size if level == 0 else font_size - 2)
            run.font.color.rgb = TEXT if level == 0 else MUTED
            run.font.bold = level == 0 and len(text) < 80


def add_two_column_slide(slide, left_title, left_items, right_title, right_items):
    header_h = 0.42

    left_header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.15), Inches(5.8), Inches(header_h)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = CYAN
    left_header.line.fill.background()
    left_header.text_frame.text = left_title
    left_header.text_frame.paragraphs[0].runs[0].font.name = "Aptos"
    left_header.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
    left_header.text_frame.paragraphs[0].runs[0].font.bold = True
    left_header.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE

    right_header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.15), Inches(5.8), Inches(header_h)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = ACCENT
    right_header.line.fill.background()
    right_header.text_frame.text = right_title
    right_header.text_frame.paragraphs[0].runs[0].font.name = "Aptos"
    right_header.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
    right_header.text_frame.paragraphs[0].runs[0].font.bold = True
    right_header.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE

    add_bullets(slide, left_items, left=0.72, top=1.68, width=5.55, height=5.35, font_size=18)
    add_bullets(slide, right_items, left=6.92, top=1.68, width=5.55, height=5.35, font_size=18)


def add_footer(slide, text="CyberShield AI"):
    box = slide.shapes.add_textbox(Inches(10.8), Inches(7.02), Inches(2.1), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(10.8), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "CyberShield AI"
    run.font.name = "Aptos Display"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NAVY

    p2 = title_box.text_frame.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Project Presentation"
    r2.font.name = "Aptos"
    r2.font.size = Pt(22)
    r2.font.color.rgb = CYAN

    summary = slide.shapes.add_textbox(Inches(0.72), Inches(2.35), Inches(11.5), Inches(1.8))
    tf = summary.text_frame
    p3 = tf.paragraphs[0]
    rr = p3.add_run()
    rr.text = (
        "An AI-powered browser-based cybersecurity dashboard that helps users detect scams, "
        "analyze links, scan files, evaluate apps, inspect public footprint exposure, and respond to incidents."
    )
    rr.font.name = "Aptos"
    rr.font.size = Pt(20)
    rr.font.color.rgb = TEXT

    cards = [
        ("Frontend", "HTML, CSS, JavaScript"),
        ("Deployment", "Netlify static hosting"),
        ("Approach", "Heuristics, scoring, NLP, APIs"),
    ]
    x_positions = [0.8, 4.55, 8.3]
    for idx, (head, body) in enumerate(cards):
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(4.6), Inches(3.1), Inches(1.35)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = RGBColor(215, 223, 232)
        t = rect.text_frame
        t.clear()
        p_head = t.paragraphs[0]
        p_head.text = head
        p_head.runs[0].font.name = "Aptos"
        p_head.runs[0].font.size = Pt(18)
        p_head.runs[0].font.bold = True
        p_head.runs[0].font.color.rgb = NAVY
        p_body = t.add_paragraph()
        p_body.text = body
        p_body.runs[0].font.name = "Aptos"
        p_body.runs[0].font.size = Pt(13)
        p_body.runs[0].font.color.rgb = MUTED

    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Problem Statement And Objective", "Why this project was built")
    add_bullets(slide, [
        "Users face many cyber risks in daily life: phishing, weak passwords, unsafe files, risky apps, fraud payments, and public privacy exposure.",
        "Most users do not have one place to check these risks quickly and understand the result.",
        "CyberShield AI was designed as a unified dashboard that provides practical security checks in one browser-based platform.",
        "The objective is to make cybersecurity support accessible, fast, and understandable for non-technical users."
    ], font_size=20)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "System Overview", "Main tools included in the dashboard")
    add_two_column_slide(
        slide,
        "Core Security Tools",
        [
            "AI Scam Detection",
            "Password Checker",
            "Real-Time File Scanner",
            "Safe Browsing Checker",
            "App Privacy Scanner",
        ],
        "Monitoring And Support",
        [
            "Transaction Risk Analyzer",
            "Digital Footprint Tracker",
            "Live Network Monitoring",
            "Live GPS Tracker",
            "Emergency Response Mode and AI Assistant",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "How The Project Works", "Execution flow from input to decision")
    add_bullets(slide, [
        "Step 1: The user selects a tool from the dashboard.",
        "Step 2: The user submits input such as message text, URL, app name, password, file, transaction details, username, phone number, or full name.",
        "Step 3: JavaScript processes the input in the browser and applies rules, scoring logic, and optional API checks.",
        "Step 4: The system generates a score, classification, explanations, red flags, and recommendations.",
        "Step 5: Results are shown instantly in the interface and can guide the user's next action.",
        "Step 6: Some tools use browser APIs or external services for richer analysis when available."
    ], font_size=19)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Programming Languages And Technologies", "Technical stack used in the project")
    add_two_column_slide(
        slide,
        "Languages",
        [
            "HTML for layout and section structure",
            "CSS for styling, layout, cards, badges, and responsive design",
            "JavaScript for all tool logic, scoring, APIs, DOM updates, and browser interactivity",
            "PowerShell used during local scripting and deployment",
        ],
        "Technologies",
        [
            "Netlify for static deployment",
            "Browser APIs: fetch, localStorage, crypto.subtle, geolocation, clipboard, DOMParser, WebSocket",
            "Optional external APIs for richer checks such as VirusTotal, GitHub, GitLab, Reddit, OpenRouter, OpenAI, and Gemini",
            "Single-page frontend architecture with modular JavaScript functions",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Algorithms Used", "Main analysis methods inside the project")
    add_bullets(slide, [
        "Rule-based heuristics: used for phishing signals, suspicious domains, risky file extensions, and fraud indicators.",
        "Weighted scoring models: used for password strength, scam risk, app privacy risk, transaction risk, and footprint confidence.",
        "Lexicon-based NLP classifier: tokenizes message text, applies suspicious and legitimate word weights, and computes phishing probability.",
        "Sigmoid-style probability conversion: converts phishing raw score into a confidence percentage.",
        "SHA-256 hashing: used in the file scanner to fingerprint uploaded files.",
        "Domain reputation logic: checks HTTPS, suspicious TLDs, IP-based URLs, subdomains, and phishing keywords.",
        "Haversine distance formula: used in GPS tracking to calculate route and nearest-location distance.",
        "Confidence ranking: sorts footprint-tracker results so stronger evidence appears first."
    ], font_size=18)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Example: AI Scam Detection", "How one major module works internally")
    add_bullets(slide, [
        "Input: pasted message, email body, or URL.",
        "Text analysis checks urgency, threats, impersonation, sensitive-data requests, reward promises, and malware lures.",
        "URL analysis checks domain tricks, suspicious TLDs, raw IP addresses, misleading brand references, and redirect behavior.",
        "NLP layer applies token weights from suspicious and legitimate vocabularies to estimate phishing probability.",
        "Legitimacy rules reduce false positives for real OTP alerts, bank notifications, academic notices, and statement emails.",
        "Output: SAFE, SUSPICIOUS, or MALICIOUS classification with confidence, risk level, reasons, and structured JSON."
    ], font_size=18)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Example: Other Analytical Modules", "How the rest of the dashboard applies security logic")
    add_two_column_slide(
        slide,
        "Security Analysis",
        [
            "Password Checker: scores passwords by complexity and common-pattern detection.",
            "File Scanner: combines suspicious extension checks, filename heuristics, SHA-256 hashing, and optional VirusTotal lookups.",
            "Safe Browsing: checks HTTPS, domain length, TLD, subdomains, raw IP use, and phishing keywords.",
            "App Privacy Scanner: aggregates trackers, permissions, reputation signals, and publisher trust.",
        ],
        "Decision Support",
        [
            "Transaction Analyzer: scores transaction context, recipient trust, amount anomaly, device trust, and payment scenario.",
            "Footprint Tracker: searches public traces for usernames, phone numbers, and full names with confidence labels.",
            "Emergency Mode: gives structured recovery steps for different compromise scenarios.",
            "AI Assistant: answers project-related and cybersecurity questions using local knowledge plus optional LLM backends.",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Real-World Implementation", "How this project is useful outside the classroom")
    add_bullets(slide, [
        "The project is implemented as a lightweight browser-based cybersecurity toolkit.",
        "Because it is a static frontend app, it is easy to host, easy to scale, and fast to access from any device.",
        "It is practical for awareness, first-level risk analysis, and guided decision support.",
        "Real users can use it to verify suspicious links, check scam messages, review password safety, inspect risky apps, and follow recovery steps after compromise.",
        "This design is realistic for cyber awareness portals, college projects, demo systems, and first-response security dashboards."
    ], font_size=19)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Why These Tools Are Required", "Security value of each module")
    add_bullets(slide, [
        "Scam Detection is required because phishing and social engineering are among the most common cyber attacks.",
        "Password tools are required because weak passwords directly lead to account compromise.",
        "File scanning is required because malicious attachments and downloads are a common attack vector.",
        "Safe browsing is required because users often click suspicious websites without checking them first.",
        "App privacy scanning is required because many mobile apps over-collect data and expose users to tracking.",
        "Transaction analysis is required because digital payment fraud is increasing rapidly.",
        "Footprint tracking is required because public exposure of names, usernames, and phone numbers creates privacy risk.",
        "Emergency response guidance is required because users need immediate action steps after a cyber incident."
    ], font_size=18)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "What Was Used To Build This Code", "Development sources, APIs, and references")
    add_two_column_slide(
        slide,
        "Development Inputs",
        [
            "Project files: HTML, CSS, and JavaScript in a static frontend structure",
            "Browser APIs: fetch, crypto.subtle, localStorage, geolocation, clipboard, DOMParser, WebSocket",
            "Netlify CLI and Netlify hosting for deployment",
            "PowerShell for local workflow and integration steps",
        ],
        "External References",
        [
            "VirusTotal API for file reputation analysis",
            "GitHub, GitLab, and Reddit public endpoints for footprint checks",
            "Google DNS over HTTPS and Gravatar for supporting lookups",
            "OpenRouter, OpenAI, and Gemini for optional AI assistant backends",
            "Adapted ideas from PHISHING-DETECTION-WITH-NLP and Cyberlayers digital_footprint_tracker",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Conclusion", "Final takeaway")
    add_bullets(slide, [
        "CyberShield AI is a multi-tool cybersecurity dashboard built to make security analysis more accessible.",
        "It combines frontend engineering, heuristic detection, scoring models, lightweight NLP, hashing, confidence ranking, and public API integrations.",
        "The project is useful because it translates complex cyber checks into understandable results for ordinary users.",
        "Its biggest strength is integration: multiple cyber safety tools in one platform with practical outputs and guided action."
    ], font_size=20)
    add_footer(slide, "Prepared From Current CyberShield AI Build")

    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
